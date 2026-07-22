"""Generate presentation PPTX for the ionization-charge unfolding paper.

Source code is the authoritative reference for the algorithm section.
Key corrections vs slides.md / JINST.tex:
  - Two-pass burst processing is the actual algorithm (V1/V3).
    The "fractional FFT phase shift" alignment from slides.md slide 7
    is NOT implemented; only global binning is used.
  - V3 processor introduces separate collection vs induction templates,
    selected per merged-group based on accumulated charge vs threshold.
  - Dead-time compensation: linear-slope interpolation, not "scaling to one tick".
  - Template compensation target is trigger_time_idx (not waveform start).
  - Active voxel threshold is set at 500 e-, matching front-end noise sigma.
  - "Hybrid mode" template option from JINST is not a named code mode;
    the code exposes center / collection / collection_plus_neighbors.
  - Wiener-inspired voxel filter is ongoing work, not yet implemented.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── palette ──────────────────────────────────────────────────────────────────
BNL_BLUE   = RGBColor(0x00, 0x47, 0x8A)   # BNL brand blue
ACCENT     = RGBColor(0xE8, 0x6C, 0x0A)   # orange accent
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0xBB, 0xBB, 0xBB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GREEN      = RGBColor(0x21, 0x7A, 0x3C)

W  = Inches(13.33)   # widescreen width
H  = Inches(7.5)     # widescreen height

# ── helpers ───────────────────────────────────────────────────────────────────

def add_slide(prs, layout_index=6):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text, l, t, w, h,
            size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
            wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE, l, t, w, h)
    # pptx uses add_shape with MSO_SHAPE_TYPE — use workaround:
    return _rect(slide, l, t, w, h, fill_color, line_color, line_width)


def _rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(1.5)):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_bullet_frame(slide, items, l, t, w, h,
                     size=16, color=DARK, bullet="  •  "):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txBox


def title_bar(slide, title, subtitle=None):
    _rect(slide, 0, 0, W, Inches(1.3), BNL_BLUE)
    textbox(slide, title,
            Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.85),
            size=28, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, subtitle,
                Inches(0.4), Inches(0.92), Inches(12.0), Inches(0.4),
                size=15, color=RGBColor(0xCC, 0xDD, 0xFF))


def section_banner(slide, label):
    _rect(slide, 0, Inches(1.3), W, Inches(0.35), ACCENT)
    textbox(slide, label,
            Inches(0.4), Inches(1.30), Inches(12.5), Inches(0.35),
            size=13, bold=True, color=WHITE)


def flow_box(slide, label, l, t, w, h,
             fill=BNL_BLUE, text_color=WHITE, size=14):
    _rect(slide, l, t, w, h, fill, line_color=None)
    textbox(slide, label, l, t, w, h,
            size=size, bold=True, color=text_color,
            align=PP_ALIGN.CENTER, wrap=True)


def arrow_down(slide, cx, y_top, length=Inches(0.28)):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    # draw a thin rectangle as arrow body
    aw = Inches(0.04)
    _rect(slide, cx - aw/2, y_top, aw, length, DARK)
    # arrowhead: tiny triangle via thin wide rect
    _rect(slide, cx - Inches(0.10), y_top + length - Inches(0.02),
          Inches(0.20), Inches(0.08), DARK)


def arrow_right(slide, x_left, cy, length=Inches(0.35)):
    ah = Inches(0.04)
    _rect(slide, x_left, cy - ah/2, length, ah, DARK)
    _rect(slide, x_left + length - Inches(0.02), cy - Inches(0.08),
          Inches(0.08), Inches(0.16), DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── 1  TITLE ──────────────────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, LIGHT_GRAY)
_rect(sl, 0, 0, W, Inches(0.18), ACCENT)
_rect(sl, 0, H - Inches(0.18), W, Inches(0.18), ACCENT)
_rect(sl, 0, Inches(1.8), W, Inches(3.5), BNL_BLUE)

textbox(sl,
        "Ionization Charge Unfolding at Zero-Suppressed\nLiquid Argon Time Projection Chamber",
        Inches(0.5), Inches(1.9), Inches(12.3), Inches(2.4),
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(sl, "An Author   |   Brookhaven National Laboratory",
        Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
        size=18, color=BNL_BLUE, align=PP_ALIGN.CENTER)

textbox(sl, "JINST paper proposal  •  Signal Processing for ND-LAr",
        Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.45),
        size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── 2  PAPER ROADMAP ──────────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Paper at a Glance", "Structure overview")

sections = [
    ("1. Introduction",         "ND-LAr physics motivation, zero-suppression challenge"),
    ("2. Signal Processing",    "Traditional deconvolution review, readout logic, algorithm"),
    ("3. Performance",          "Simulation setup, template/filter/threshold/noise studies"),
    ("4. Discussion",           "Truth bookkeeping, hardware, readout requirements, shield grid"),
    ("5. Summary",              "Conclusions and outlook"),
]
y0 = Inches(1.55)
dy = Inches(0.98)
for i, (sec, desc) in enumerate(sections):
    y = y0 + i * dy
    _rect(sl, Inches(0.35), y, Inches(3.2), Inches(0.82), BNL_BLUE)
    textbox(sl, sec, Inches(0.35), y, Inches(3.2), Inches(0.82),
            size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, desc, Inches(3.75), y + Inches(0.15), Inches(9.0), Inches(0.6),
            size=14, color=DARK)

# ── 3  MOTIVATION ─────────────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Motivation", "Why zero suppression breaks conventional signal processing")
section_banner(sl, "Introduction")

col_l = Inches(0.35)
col_r = Inches(6.9)
cw    = Inches(6.1)

textbox(sl, "ND-LAr facts",
        col_l, Inches(1.8), cw, Inches(0.35),
        size=16, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Pixelated readout  (LArPix-v2 ASICs)",
    "Self-triggering: electronics dormant until charge exceeds discriminator threshold",
    "ADC fires only after a configurable hold delay (adc_hold_delay)",
    "Front-end resets and drains charge after each burst sequence",
    "Burst mode: records several consecutive integration windows without reset",
], col_l, Inches(2.15), cw, Inches(2.4), size=14)

textbox(sl, "Consequences for signal processing",
        col_r, Inches(1.8), cw, Inches(0.35),
        size=16, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Waveform is truncated: only charge above threshold is recorded",
    "First trigger time is known; subsequent hits lose exact start time",
    "Negative induction signals from neighboring pixels are never recorded",
    "Reset mechanism may cause charge to be lost or overestimated",
    "Traditional deconvolution assumes full, equally-spaced waveform sampling",
], col_r, Inches(2.15), cw, Inches(2.4), size=14)

_rect(sl, Inches(0.35), Inches(4.75), W - Inches(0.7), Inches(1.0), LIGHT_GRAY)
textbox(sl,
        "Goal:  Recover approximate full waveforms using physics-guided templates, "
        "then apply FFT-based 3D deconvolution to unfold the ionization charge.",
        Inches(0.55), Inches(4.82), W - Inches(1.1), Inches(0.88),
        size=15, bold=True, color=BNL_BLUE)

# ── 4  TRADITIONAL SIGNAL PROCESSING ─────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Traditional Signal Processing", "What works at the far detector")
section_banner(sl, "Section 2 review")

textbox(sl, "System model:   m(t) = r(t) * s(t) + n(t)",
        Inches(0.4), Inches(1.75), Inches(8.0), Inches(0.5),
        size=17, bold=True, color=BNL_BLUE)

add_bullet_frame(sl, [
    "m(t)  measured waveform,  r(t)  detector + electronics response,  s(t)  ionization signal,  n(t)  noise",
    "Deconvolution in frequency domain:  S(w) = M(w) / R(w)  (minus noise term)",
    "Gaussian filter suppresses blow-up at small R(w) while preserving the charge integral",
    "Region-of-interest (ROI) selection removes signal-free time bins",
    "Wiener-inspired filters give simultaneous spatial and charge resolution",
], Inches(0.4), Inches(2.25), Inches(12.5), Inches(2.4), size=14)

textbox(sl, "Why this fails at ND-LAr",
        Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.35),
        size=16, bold=True, color=ACCENT)
add_bullet_frame(sl, [
    "Zero suppression destroys the continuous, equally-spaced waveform required by FFT",
    "The recorded signal is no longer a linear function of ionization charge",
    "Direct convolution / deconvolution on raw measurements is invalid",
], Inches(0.4), Inches(5.05), Inches(12.5), Inches(1.4), size=14)

# ── 5  READOUT LOGIC ──────────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "ND-LAr Readout Logic", "Self-triggering and burst mode details")
section_banner(sl, "Section 2 — readout review")

col_l = Inches(0.35)
col_r = Inches(6.9)
cw    = Inches(6.1)

textbox(sl, "Self-trigger cycle (single hit)",
        col_l, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Charge accumulates on CSA front-end",
    "Discriminator fires when charge > threshold",
    "After adc_hold_delay, ADC latches charge if still above threshold (veto random noise)",
    "During ADC conversion, no new triggers possible  =>  dead time",
    "Front-end resets; new charge accumulation starts",
    "First trigger time is recorded precisely; consecutive hits lose exact start time",
], col_l, Inches(2.1), cw, Inches(2.9), size=13)

textbox(sl, "Burst mode extension",
        col_r, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Front-end does NOT reset after the first hold window",
    "Records several more charge integrals at fixed time windows (adc_hold_delay)",
    "Resets only at the end of the burst sequence",
    "Burst length (nburst) is a configurable parameter",
    "Gives more waveform samples per trigger at the cost of saturation risk",
    "Code: t_first = trigger_time + adc_hold_delay;  t_last = trigger_time + adc_hold_delay * nburst",
], col_r, Inches(2.1), cw, Inches(2.9), size=13)

_rect(sl, Inches(0.35), Inches(5.25), W - Inches(0.7), Inches(0.7), LIGHT_GRAY)
textbox(sl,
        "Key quantity: the gap between two burst sequences  =  t_B_start - t_A_last  "
        "(difference of burst start times, not integration end times)",
        Inches(0.55), Inches(5.30), W - Inches(1.1), Inches(0.6),
        size=13, italic=True, color=DARK)

# ── 6  ALGORITHM OVERVIEW ─────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Algorithm Overview", "Three-stage pipeline")
section_banner(sl, "Section 2 — algorithm")

# Flow chart boxes
bw = Inches(3.4)
bh = Inches(0.75)
gap = Inches(0.45)
y0  = Inches(1.75)
cx  = [Inches(0.45), Inches(4.95), Inches(9.35)]
colors = [BNL_BLUE, BNL_BLUE, BNL_BLUE]
labels = [
    "Stage 1\nWaveform Recovery",
    "Stage 2\n3D FFT Deconvolution",
    "Stage 3\nActive Voxel Selection",
]
for i in range(3):
    flow_box(sl, labels[i], cx[i], y0, bw, bh, fill=colors[i], size=15)
    if i < 2:
        arrow_right(sl, cx[i] + bw, y0 + bh/2 - Inches(0.02), gap)

# Sub-steps below each stage
sub = [
    [
        "Extract burst sequences per pixel",
        "Pass 1: dead-time merge (gap <= tau)",
        "Pass 2: template compensation (gap > tau)",
        "Assemble into 3D block",
    ],
    [
        "Load & integrate field response kernel",
        "Build 3D Gaussian filter",
        "FFT deconvolution: S = M / R  x  G",
        "Trim & roll output to correct offset",
    ],
    [
        "Apply 3D Gaussian filter to deconv output",
        "Select voxels with charge > 500 e-",
        "(Ongoing) Wiener-inspired frequency filter",
        "Compare with smeared truth for residual",
    ],
]
for i, items in enumerate(sub):
    add_bullet_frame(sl, items, cx[i], y0 + bh + Inches(0.1),
                     bw, Inches(2.4), size=12, color=DARK)

textbox(sl,
        "Input: raw LArPix hit data  (pixel x/y, trigger time, burst charge integrals)     "
        "Output: 3D deconvolved charge image",
        Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
        size=12, italic=True, color=MID_GRAY)

# ── 7  STAGE 1a — BURST SEQUENCE EXTRACTION ──────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Stage 1 — Burst Sequence Extraction", "From raw hits to per-pixel sequences")
section_banner(sl, "Algorithm detail")

textbox(sl, "Raw hit data format (Hits container)",
        Inches(0.35), Inches(1.75), Inches(6.0), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
_rect(sl, Inches(0.35), Inches(2.1), Inches(6.0), Inches(1.6), LIGHT_GRAY)
textbox(sl,
        "location[:, 0]  pixel_x\n"
        "location[:, 1]  pixel_y\n"
        "location[:, 2]  trigger_time_idx\n"
        "location[:, 3]  last_adc_latch\n"
        "location[:, 4]  next_integration_start\n"
        "data[:, 3:]     cumulative charge per burst window",
        Inches(0.5), Inches(2.15), Inches(5.7), Inches(1.5),
        size=13, color=DARK)

textbox(sl, "Per-pixel BurstSequence construction",
        Inches(6.6), Inches(1.75), Inches(6.3), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "charges = diff(cumulative data)  =>  charge per burst window",
    "t_first = trigger_time_idx + adc_hold_delay",
    "t_last  = trigger_time_idx + adc_hold_delay * nburst",
    "Sequences grouped by (pixel_x, pixel_y) and sorted by trigger_time_idx",
    "Validation: no overlapping sequences (t_first[k+1] >= t_last[k])",
], Inches(6.6), Inches(2.1), Inches(6.3), Inches(2.4), size=13)

textbox(sl, "Gap definition",
        Inches(0.35), Inches(3.85), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
_rect(sl, Inches(0.35), Inches(4.2), W - Inches(0.7), Inches(0.65), LIGHT_GRAY)
textbox(sl,
        "gap = seq_B.t_first - seq_A.t_last     "
        "(gap between the start of B and the start of the last burst of A, "
        "NOT the end of A's integration window)",
        Inches(0.55), Inches(4.25), W - Inches(1.1), Inches(0.55),
        size=13, italic=True, color=DARK)

add_bullet_frame(sl, [
    "gap = 0:  sequences are adjacent (touching) but not overlapping",
    "0 < gap <= tau:  close enough for dead-time compensation  (Pass 1)",
    "gap > tau:  template compensation required  (Pass 2)",
], Inches(0.35), Inches(5.0), Inches(12.5), Inches(1.3), size=14)

# ── 8  STAGE 1b — TWO-PASS WAVEFORM RECOVERY ─────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Stage 1 — Two-Pass Waveform Recovery", "Dead-time merge then template fill")
section_banner(sl, "Algorithm detail  (BurstSequenceProcessorV3)")

# Pass 1 box
_rect(sl, Inches(0.35), Inches(1.75), Inches(5.9), Inches(2.55), BNL_BLUE)
textbox(sl, "Pass 1 — Dead-Time Compensation",
        Inches(0.45), Inches(1.77), Inches(5.7), Inches(0.45),
        size=15, bold=True, color=WHITE)
add_bullet_frame(sl, [
    "Iterate over sequences for a pixel in time order",
    "If gap <= tau:  merge seq_A and seq_B",
    "   slope = charge_B[0] / (gap - deadtime)",
    "   compensated = slope * deadtime",
    "   first charge of B += compensated",
    "Output: list of MergedBurstGroups (one per cluster)",
], Inches(0.45), Inches(2.2), Inches(5.7), Inches(2.0), size=12, color=WHITE)

arrow_right(sl, Inches(6.25), Inches(3.05), Inches(0.55))

# Pass 2 box
_rect(sl, Inches(6.8), Inches(1.75), Inches(6.15), Inches(2.55), ACCENT)
textbox(sl, "Pass 2 — Template Compensation",
        Inches(6.9), Inches(1.77), Inches(5.95), Inches(0.45),
        size=15, bold=True, color=WHITE)
add_bullet_frame(sl, [
    "Treat each MergedBurstGroup as a block",
    "Select template: collection (max cumq > threshold) or",
    "   induction positive prefix (max cumq <= threshold)",
    "Find transit fraction: threshold / max_cumulative_charge",
    "Scan template to find window of length = gap - deadtime",
    "   that rises by transit fraction",
    "Interpolate template points, append before group",
    "Scale so template endpoint equals threshold",
    "Output: continuous MergedSequence per pixel",
], Inches(6.9), Inches(2.2), Inches(5.9), Inches(2.05), size=12, color=WHITE)

# Bootstrap note
_rect(sl, Inches(0.35), Inches(4.45), W - Inches(0.7), Inches(0.65), LIGHT_GRAY)
textbox(sl,
        "Bootstrap (first sequence):  no previous time; scan template from its start to find "
        "where cumulative equals transit fraction, then work backwards to trigger_time_idx.",
        Inches(0.55), Inches(4.5), W - Inches(1.1), Inches(0.55),
        size=12, italic=True, color=DARK)

textbox(sl, "Block assembly",
        Inches(0.35), Inches(5.25), Inches(3.0), Inches(0.35),
        size=14, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "All MergedSequences binned into 3D array (pixel_x, pixel_y, time_bin)",
    "Time bin size = adc_hold_delay;  padded by npadbin bins on each side",
], Inches(0.35), Inches(5.6), Inches(12.5), Inches(0.9), size=13)

# ── 9  BURST PROCESSOR VERSIONS ──────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Burst Processor Versions", "V1 vs V2 vs V3 strategies")
section_banner(sl, "Algorithm variants  (burst_processor.py / v2.py / v3.py)")

# Table-like comparison
y0 = Inches(1.75)
cols = [
    ("Aspect", Inches(0.35), Inches(2.3)),
    ("V1", Inches(2.8), Inches(1.8)),
    ("V2", Inches(4.7), Inches(1.8)),
    ("V3", Inches(6.6), Inches(2.4)),
]

header_y = y0
for label, x, w in cols:
    _rect(sl, x, header_y, w, Inches(0.45), BNL_BLUE)
    textbox(sl, label, x, header_y, w, Inches(0.45),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Passes", "Single pass", "Single pass", "Two passes"),
    ("Dead-time merge", "Slope compensation", "Slope compensation", "Pass 1 only"),
    ("Time alignment", "Global binning only", "Fractional FFT phase shift", "Global binning only"),
    ("Template mode", "Single (cumulative)", "Single (cumulative)", "Dual: collection + induction"),
    ("Template selection", "Fixed per pixel", "Fixed per pixel", "Per merged group by charge"),
]

y = header_y + Inches(0.45)
for row_label, v1_val, v2_val, v3_val in rows:
    # Label column
    _rect(sl, Inches(0.35), y, Inches(2.3), Inches(0.55), LIGHT_GRAY)
    textbox(sl, row_label, Inches(0.35), y, Inches(2.3), Inches(0.55),
            size=11, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    # V1
    _rect(sl, Inches(2.8), y, Inches(1.8), Inches(0.55), WHITE)
    textbox(sl, v1_val, Inches(2.8), y, Inches(1.8), Inches(0.55),
            size=10, color=DARK, align=PP_ALIGN.CENTER)

    # V2
    _rect(sl, Inches(4.7), y, Inches(1.8), Inches(0.55), WHITE)
    textbox(sl, v2_val, Inches(4.7), y, Inches(1.8), Inches(0.55),
            size=10, color=DARK, align=PP_ALIGN.CENTER)

    # V3
    _rect(sl, Inches(6.6), y, Inches(2.4), Inches(0.55), WHITE)
    textbox(sl, v3_val, Inches(6.6), y, Inches(2.4), Inches(0.55),
            size=10, color=DARK, align=PP_ALIGN.CENTER)

    y += Inches(0.55)

# Notes below
textbox(sl, "Notes:",
        Inches(0.35), Inches(5.8), Inches(12.5), Inches(0.3),
        size=14, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "V1: baseline single-pass processor; most direct implementation of dead-time + template compensation.",
    "V2: adds sub-sample time-jitter correction via frequency-domain phase shifts; higher complexity; not used in current pipeline.",
    "V3 (current): two-pass strategy separates close-sequence merging from distant-gap filling; introduces charge-dependent template selection.",
], Inches(0.35), Inches(6.15), Inches(12.5), Inches(1.2), size=12)

# ── 9b STAGE 1 — TEMPLATE DETAILS ────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Stage 1 — Field Response Templates", "What gets used as the physics template")
section_banner(sl, "Algorithm detail  (deconv_workflow.py  /  field_response.py)")

textbox(sl, "Three template modes",
        Inches(0.35), Inches(1.75), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)

modes = [
    ("center",
     "Response at the center of the collection pixel only.\n"
     "Fast, symmetric, good approximation for high-charge pixels."),
    ("collection",
     "Average response over all intra-pixel sampling paths of the collection pixel.\n"
     "Reduces path-dependent smearing."),
    ("collection_plus_neighbors",
     "Average over collection pixel plus the 8 surrounding neighbors (radius=1).\n"
     "Captures long-range induction; uses positive_cumulative template search mode.\n"
     "V3 processor further splits: collection template for pixels with max cumq > threshold,\n"
     "induction (positive-prefix only) for pixels with max cumq <= threshold."),
]
y = Inches(2.15)
for name, desc in modes:
    _rect(sl, Inches(0.35), y, Inches(2.8), Inches(1.15), BNL_BLUE)
    textbox(sl, name, Inches(0.35), y, Inches(2.8), Inches(1.15),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, desc, Inches(3.3), y + Inches(0.1), Inches(9.7), Inches(1.0),
            size=13, color=DARK)
    y += Inches(1.25)

textbox(sl, "Template preprocessing",
        Inches(0.35), Inches(5.55), Inches(6.0), Inches(0.35),
        size=14, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Raw field response multiplied by time_tick, spatial quadrant mirrored to full plane",
    "Kernel flipped for convolution; intra-pixel paths downsampled by averaging",
    "Cumulative template built from cumsum, then maximum.accumulate enforces monotone increase  (handles bipolar dips)",
    "Induction template for V3: keep only leading positive prefix of bipolar waveform",
], Inches(0.35), Inches(5.9), Inches(12.5), Inches(1.3), size=13)

# ── 10  STAGE 2 — 3D FFT DECONVOLUTION ───────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Stage 2 — 3D FFT Deconvolution", "Inverting the detector response in frequency space")
section_banner(sl, "Algorithm detail  (deconv.py)")

col_l = Inches(0.35)
col_r = Inches(6.9)
cw    = Inches(6.1)

textbox(sl, "Inputs",
        col_l, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Recovered 3D charge block  (pixel_x, pixel_y, time_bin)",
    "Field response kernel: averaged over intra-pixel paths, then summed per adc_hold_delay time window",
    "   kernel shape:  (n_pix_x, n_pix_y, n_time_bins)",
    "3D Gaussian filter   sigma_pixel x sigma_pixel x sigma_time",
], col_l, Inches(2.1), cw, Inches(1.8), size=13)

textbox(sl, "Deconvolution steps",
        col_r, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Zero-pad both block and kernel to extended shape  (avoids circular aliasing)",
    "rfftn(block)  and  rfftn(kernel)  =>  frequency-domain arrays",
    "Divide element-wise:  S_hat = M_hat / R_hat",
    "Multiply by Gaussian filter  G_hat  to suppress high-frequency noise",
    "irfftn  =>  real-space deconvolved charge",
    "Roll by kernel half-width to re-center output",
    "Trim to expected spatial and time extent",
], col_r, Inches(2.1), cw, Inches(2.5), size=13)

textbox(sl, "3D Gaussian filter",
        col_l, Inches(4.1), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Separable product of 1D Gaussians along each axis",
    "Pixel axes:  sigma_pixel  (in units of pixel pitch)",
    "Time axis:   sigma_time   (in same units as adc_hold_delay)",
    "Built in frequency space via rfftfreq / fftfreq for each axis",
    "Typical values:  sigma_pixel ~ 0.8 pitch,  sigma_time ~ 1.6 to several us",
], col_l, Inches(4.45), cw, Inches(2.0), size=13)

textbox(sl, "Complexity",
        col_r, Inches(4.1), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "FFT-based:  O(N log N)  where N = total number of voxels",
    "Compatible with GPU acceleration via PyTorch / CuPy rfftn",
    "Block-wise deconvolution planned to reduce memory footprint",
], col_r, Inches(4.45), cw, Inches(1.3), size=13)

# ── 11  STAGE 3 — ACTIVE VOXEL SELECTION ─────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Stage 3 — Active Voxel Identification", "Selecting reconstructed signal voxels")
section_banner(sl, "Algorithm detail")

textbox(sl, "Current strategy",
        Inches(0.35), Inches(1.75), Inches(6.0), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Gaussian filter applied to deconvolved 3D array",
    "Voxels with filtered charge > 500 e- are kept as active",
    "500 e- corresponds approximately to the LArPix-v2 front-end noise standard deviation",
], Inches(0.35), Inches(2.1), Inches(6.0), Inches(1.5), size=14)

textbox(sl, "Truth comparison via Gaussian smearing",
        Inches(6.6), Inches(1.75), Inches(6.3), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "The same 3D Gaussian filter is applied to the true ionization charge (from simulation)",
    "Active deconvolved voxels define a set of (pixel, time-bin) pairs",
    "True charge in matching pairs is summed and linked to particle ID",
    "Residual = deconvolved charge - smeared true charge  per voxel",
    "This allows fair comparison without mismatch from detector smearing",
], Inches(6.6), Inches(2.1), Inches(6.3), Inches(2.3), size=13)

textbox(sl, "Ongoing: Wiener-inspired frequency filter for voxel selection",
        Inches(0.35), Inches(3.75), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=ACCENT)
add_bullet_frame(sl, [
    "Form:  F(w) = exp( -(w/a)^b )  applied in frequency domain",
    "Goal: adaptively suppress noise-dominated frequency bands",
    "Will replace fixed-threshold selection with a signal-to-noise driven criterion",
    "Under active development; not yet in the main pipeline",
], Inches(0.35), Inches(4.1), Inches(12.5), Inches(1.5), size=14)

_rect(sl, Inches(0.35), Inches(5.8), W - Inches(0.7), Inches(0.7), LIGHT_GRAY)
textbox(sl,
        "Output payload (NPZ):  deconv_q, boffset, smeared_true, smear_offset, "
        "hits data, effective-charge truth, template-compensation diagnostics, geometry metadata",
        Inches(0.55), Inches(5.85), W - Inches(1.1), Inches(0.6),
        size=12, italic=True, color=DARK)

# ── 12  SIMULATION SETUP ──────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Simulation Setup", "How performance is evaluated")
section_banner(sl, "Section 3 — performance")

col_l = Inches(0.35)
col_r = Inches(6.9)
cw    = Inches(6.1)

textbox(sl, "Detector and input",
        col_l, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Geometry:  ProtoDUNE-ND  (2x2 experiment modules)",
    "Input:     positron particle gun",
    "Simulator: TRED  (GPU-accelerated, provides true ionization charge with particle ID)",
    "Provides: hit data, effective charge truth (effq), waveforms",
], col_l, Inches(2.1), cw, Inches(2.0), size=14)

textbox(sl, "Metrics",
        col_r, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Charge residual:  deconvolved charge minus smeared true charge per active voxel",
    "Residual width characterizes charge reconstruction resolution",
    "Studies vary: template choice, filter width, threshold, noise level",
], col_r, Inches(2.1), cw, Inches(2.0), size=14)

textbox(sl, "Study matrix",
        col_l, Inches(4.3), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)

studies = [
    "Template",
    "Filter",
    "Threshold",
    "Noise",
]
descs = [
    "center / collection / collection+neighbors / (ongoing: filtered template)",
    "Pixel Gaussian sigma ~ 0.8 pitch;  time sigma from 1.6 us to several us",
    "No threshold vs 1000 e- vs 5000 e- (lower threshold preserves more waveform info)",
    "With and without simulated LArPix-v2 noise  (noise is subleading effect)",
]
y = Inches(4.65)
for s, d in zip(studies, descs):
    _rect(sl, Inches(0.35), y, Inches(2.2), Inches(0.48), BNL_BLUE)
    textbox(sl, s, Inches(0.35), y, Inches(2.2), Inches(0.48),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, d, Inches(2.7), y + Inches(0.07), Inches(10.3), Inches(0.38),
            size=13, color=DARK)
    y += Inches(0.58)

# ── 13  PERFORMANCE SUMMARY ───────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Performance Results", "What the numbers show")
section_banner(sl, "Section 3 — performance")

_rect(sl, Inches(0.35), Inches(1.75), W - Inches(0.7), Inches(1.05), BNL_BLUE)
textbox(sl,
        "Key result:  charge residual width < 500 electrons across evaluated configurations",
        Inches(0.55), Inches(1.82), W - Inches(1.1), Inches(0.9),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

findings = [
    ("Templates",
     "Including neighboring-pixel induction (collection_plus_neighbors) improves recovery "
     "and reduces ghost hits.  Center-only template overestimates charge on induction pixels."),
    ("Filters",
     "Pixel-axis Gaussian at sigma ~ 0.8 pitch provides good spatial resolution. "
     "Temporal sigma of 1.6 us (one digitization period) preserves timing while suppressing noise."),
    ("Threshold",
     "Lower threshold preserves more of the rising waveform and reduces start-time ambiguity. "
     "5000 e- threshold significantly degrades deconvolution quality."),
    ("Noise",
     "Noise is a subleading contribution. Waveform truncation and threshold effects "
     "dominate reconstruction uncertainty in the current setup."),
]
y = Inches(2.95)
for title, desc in findings:
    _rect(sl, Inches(0.35), y, Inches(2.4), Inches(0.75), ACCENT)
    textbox(sl, title, Inches(0.35), y, Inches(2.4), Inches(0.75),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, desc, Inches(2.95), y + Inches(0.08), Inches(10.0), Inches(0.65),
            size=13, color=DARK)
    y += Inches(0.9)

# ── 14  DISCUSSION: TRUTH BOOKKEEPING ────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Discussion — Truth Bookkeeping", "Linking reconstructed charge to simulation truth")
section_banner(sl, "Section 4 — discussion")

textbox(sl, "Problem",
        Inches(0.35), Inches(1.75), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
textbox(sl,
        "True ionization charge and deconvolved charge cannot be compared directly "
        "because detector smearing affects both.  A naive point-by-point comparison "
        "would suffer from spatial mismatch.",
        Inches(0.35), Inches(2.1), Inches(12.5), Inches(0.7),
        size=14, color=DARK)

textbox(sl, "Solution: apply the same Gaussian smearing to both sides",
        Inches(0.35), Inches(2.95), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "TRED saves true ionization charge with particle ID per voxel",
    "Apply identical 3D Gaussian filter (sigma_pixel, sigma_time) to the true charge",
    "Define active voxel set from deconvolved output",
    "Sum smeared true charge within each active (pixel, time-bin) pair",
    "Link total true charge in each pair to the particle ID with the largest contribution",
    "Residual = deconvolved charge - smeared true charge  per active voxel",
], Inches(0.35), Inches(3.3), Inches(12.5), Inches(2.2), size=14)

_rect(sl, Inches(0.35), Inches(5.65), W - Inches(0.7), Inches(0.72), LIGHT_GRAY)
textbox(sl,
        "This provides a direct bridge between the signal-processing chain and "
        "machine-learning-based reconstruction workflows that rely on ionization history tracking.",
        Inches(0.55), Inches(5.7), W - Inches(1.1), Inches(0.62),
        size=13, italic=True, color=DARK)

# ── 15  DISCUSSION: HARDWARE & READOUT ───────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Discussion — Hardware & Readout Requirements", "What the detector needs to do well")
section_banner(sl, "Section 4 — discussion")

col_l = Inches(0.35)
col_r = Inches(6.9)
cw    = Inches(6.1)

textbox(sl, "GPU acceleration (ongoing benchmark)",
        col_l, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "FFT deconvolution is O(N log N); scales favorably with detector size",
    "PyTorch / CuPy rfftn enables GPU offload with minimal code change",
    "CPU vs GPU benchmark is planned but not yet complete",
], col_l, Inches(2.1), cw, Inches(1.5), size=13)

textbox(sl, "Readout design recommendations",
        col_r, Inches(1.75), cw, Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Low threshold:  preserves rising waveform, reduces start-time ambiguity",
    "Short burst sequence:  limits missing waveform fraction",
    "Threshold uniformity across pixels is important (todo: quantify)",
    "Saturation: non-reset during burst means charge can saturate at high topology density",
], col_r, Inches(2.1), cw, Inches(1.85), size=13)

textbox(sl, "Shield grid study (DUNE internal)",
        col_l, Inches(4.05), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=BNL_BLUE)
add_bullet_frame(sl, [
    "Ramo theorem: adding a shield grid sharpens the collection-pixel signal",
    "Suppresses long-range induction on neighboring pixels",
    "Shorter effective waveform  =>  smaller missing-waveform fraction",
    "Expected results: cleaner event display, fewer ghost hits, smaller residual width",
    "Demonstrated using same simulation setup with and without shield grid",
], col_l, Inches(4.4), Inches(12.5), Inches(1.8), size=14)

# ── 16  CONCLUSIONS ──────────────────────────────────────────────────────────
sl = add_slide(prs)
bg(sl, WHITE)
title_bar(sl, "Conclusions", "")

_rect(sl, 0, 0, W, Inches(0.18), ACCENT)  # already added by title_bar bg... repeat ok

contributions = [
    "Physics-guided waveform recovery via two-pass burst processing "
    "(dead-time merge + template compensation)",
    "Separate collection and induction templates for pixels with different charge levels (V3)",
    "FFT-based 3D deconvolution with integrated field-response kernel and Gaussian filter",
    "Charge residual width below 500 electrons demonstrated in simulation",
    "Truth bookkeeping framework linking reconstructed charge to simulation particle IDs",
]
y = Inches(1.7)
for c in contributions:
    _rect(sl, Inches(0.35), y, Inches(0.45), Inches(0.55), ACCENT)
    textbox(sl, str(contributions.index(c)+1),
            Inches(0.35), y, Inches(0.45), Inches(0.55),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, c, Inches(0.9), y + Inches(0.07),
            Inches(12.05), Inches(0.48), size=14, color=DARK)
    y += Inches(0.68)

textbox(sl, "Open items",
        Inches(0.35), Inches(5.4), Inches(12.5), Inches(0.35),
        size=15, bold=True, color=ACCENT)
add_bullet_frame(sl, [
    "Wiener-inspired frequency filter for voxel selection (ongoing)",
    "Block-wise deconvolution for large events (planned)",
    "CPU vs GPU benchmark (planned)",
    "Threshold uniformity and saturation studies (planned)",
    "Shield grid study results (DUNE internal review)",
], Inches(0.35), Inches(5.75), Inches(12.5), Inches(1.55), size=13)

# ─────────────────────────────────────────────────────────────────────────────
out = "/home/yousen/Documents/NDLAr2x2/sp_deconv_tradition/UnfoldLArPix/docs/slides_unfoldlarpix.pptx"
prs.save(out)
print(f"Saved: {out}")
